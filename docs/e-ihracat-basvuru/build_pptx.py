#!/usr/bin/env python3
"""Generate the OpenIDX e-ihracat application deck as an editable PPTX,
mirroring the HTML/PDF version's content and dark theme. 16:9."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# palette (matches the HTML deck)
BG       = RGBColor(0x0B, 0x12, 0x20)
PANEL    = RGBColor(0x0F, 0x17, 0x2A)
INK      = RGBColor(0xE2, 0xE8, 0xF0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
BRAND2   = RGBColor(0x60, 0xA5, 0xFA)
GOLD     = RGBColor(0xFB, 0xBF, 0x24)
OK       = RGBColor(0x34, 0xD3, 0x99)
BAD      = RGBColor(0xF8, 0x71, 0x71)
LINE     = RGBColor(0x1E, 0x29, 0x3B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.adjustments[0] = 0.06
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def txt(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    first = True
    # runs: list of paragraphs; each paragraph = list of (text,size,color,bold)
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(sp_after)
        for (text, size, color, bold) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = "Segoe UI"
    return tb


def footer(slide, n):
    txt(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
        [[("E-İhracat Başvuru Sunumu · OpenIDX", 10, RGBColor(0x47,0x55,0x69), False)]])
    txt(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.3),
        [[(str(n), 10, RGBColor(0x47,0x55,0x69), False)]], align=PP_ALIGN.RIGHT)


def kicker_title(slide, kicker, title, kcolor=BRAND2):
    txt(slide, Inches(0.65), Inches(0.5), Inches(12), Inches(0.4),
        [[(kicker.upper(), 13, kcolor, True)]])
    txt(slide, Inches(0.65), Inches(0.9), Inches(12), Inches(0.9),
        [[(title, 32, WHITE, True)]])


def kpi(slide, l, t, w, h, num, label):
    box(slide, l, t, w, h, fill=PANEL, line=LINE)
    txt(slide, l, t+Inches(0.28), w, Inches(0.9),
        [[(num, 34, BRAND2, True)]], align=PP_ALIGN.CENTER)
    txt(slide, l+Inches(0.12), t+Inches(1.25), w-Inches(0.24), h-Inches(1.3),
        [[(label, 11, MUTED, False)]], align=PP_ALIGN.CENTER)


def tile(slide, l, t, w, h, k, title, desc):
    box(slide, l, t, w, h, fill=PANEL, line=LINE)
    txt(slide, l+Inches(0.2), t+Inches(0.16), w-Inches(0.4), h-Inches(0.3),
        [[(k.upper(), 11, BRAND2, True)],
         [(title, 16, WHITE, True)],
         [(desc, 11, MUTED, False)]], sp_after=4)


def bullet_box(slide, l, t, w, h, header, hcolor, items, icolor=INK, fill=None, line=LINE):
    box(slide, l, t, w, h, fill=fill, line=line)
    paras = [[(header, 16, hcolor, True)]]
    for it in items:
        paras.append([("• ", 13, hcolor, True), (it, 13, icolor, False)])
    txt(slide, l+Inches(0.22), t+Inches(0.18), w-Inches(0.44), h-Inches(0.3),
        paras, sp_after=5)


# ---------- S1 cover ----------
s = prs.slides.add_slide(BLANK); bg(s)
txt(s, Inches(0.65), Inches(0.55), Inches(6), Inches(0.5),
    [[("◇  Open", 22, WHITE, True), ("IDX", 22, BRAND2, True)]])
txt(s, Inches(0.65), Inches(1.7), Inches(12), Inches(0.4),
    [[("E-İHRACAT DESTEK BAŞVURUSU · TANITIM SUNUMU", 14, BRAND2, True)]])
txt(s, Inches(0.65), Inches(2.2), Inches(12), Inches(1.6),
    [[("Dört güvenlik ürününü tek platformda birleştiren,", 40, WHITE, True)],
     [("self-host siber güvenlik yazılımı", 40, BRAND2, True)]], sp_after=2)
txt(s, Inches(0.65), Inches(4.3), Inches(10.5), Inches(1.0),
    [[("Kimlik (IAM) · Yönetişim (IGA) · Ayrıcalıklı Erişim (PAM) · Sıfır-Güven Ağ (ZTNA) — "
       "tek kod tabanı, tek PostgreSQL. Dünya çapında dijital olarak dağıtılan, yerli üretim bir yazılım ürünü.",
       17, MUTED, False)]])
b = box(s, Inches(0.65), Inches(5.5), Inches(6.6), Inches(0.55), fill=PANEL, line=GOLD)
txt(s, Inches(0.85), Inches(5.6), Inches(6.3), Inches(0.4),
    [[("Olgunluk: MVP–Erken GA · Ücretli pilot ve tasarım-ortağı için hazır", 12.5, GOLD, True)]])
footer(s, 1)

# ---------- S2 executive summary ----------
s = prs.slides.add_slide(BLANK); bg(s)
kicker_title(s, "Yönetici Özeti", "Bir bakışta OpenIDX")
kw = Inches(2.95); kh = Inches(1.9); ky = Inches(1.95)
kpi(s, Inches(0.65), ky, kw, kh, "4→1", "Dört ayrı ürün yerine tek platform (Okta+SailPoint+CyberArk+Zscaler)")
kpi(s, Inches(3.75), ky, kw, kh, "%70-80", "Yığılmış rakip lisans maliyetine göre tasarruf")
kpi(s, Inches(6.85), ky, kw, kh, "100%", "Self-host · veri müşterinin altyapısında")
kpi(s, Inches(9.95), ky, kw, kh, "SaaS", "Dijital teslimat · kurulum ve lisans çevrimiçi")
bullet_box(s, Inches(0.65), Inches(4.15), Inches(6.0), Inches(2.5),
           "Neden ihraç edilebilir bir üründür?", OK,
           ["Yazılım = fiziksel gümrük yok; dünyaya dijital kurulur",
            "Açık kaynak çekirdek → küresel geliştirici hunisi (GitHub)",
            "Self-host'u zorunlu kılan her regüle ülke potansiyel pazar"],
           fill=RGBColor(0x0E,0x1E,0x18))
bullet_box(s, Inches(7.0), Inches(4.15), Inches(5.65), Inches(2.5),
           "Farklılaştırıcı (moat)", BRAND2,
           ["Tek kill-switch: token+oturum+kasa+ağ devresini ≤30 sn'de keser",
            "Rakipler çok-satıcılı entegrasyonla kurarken OpenIDX tek işlemle yapar"],
           fill=PANEL)
footer(s, 2)

# ---------- S3 section 1 e-ihracat ----------
s = prs.slides.add_slide(BLANK); bg(s)
kicker_title(s, "BÖLÜM 1", "Projenin e-ihracatla ilişkisi", kcolor=GOLD)
txt(s, Inches(0.65), Inches(1.75), Inches(12), Inches(0.7),
    [[("OpenIDX bir yazılım ürünüdür; ihracatı fiziksel değil, tamamen dijitaldir. "
       "Ürün internet üzerinden teslim edilir, lisanslanır ve desteklenir.", 18, INK, False)]])
tw = Inches(6.0); th = Inches(1.55)
tile(s, Inches(0.65), Inches(2.65), tw, th, "Dijital teslimat", "Sınır ötesi, gümrüksüz",
     "Container imajları/kurulum paketleri çevrimiçi dağıtılır. Yurt dışı müşteri kendi altyapısına dakikalar içinde kurar; nakliye/gümrük/stok yok.")
tile(s, Inches(7.0), Inches(2.65), tw, th, "Gelir modeli", "Yinelenen döviz geliri",
     "Kullanıcı-başı yıllık abonelik (Enterprise), destek/SLA ve uyum eklentileri. Yurt dışı satış = döviz cinsinden yinelenen ihracat geliri.")
tile(s, Inches(0.65), Inches(4.35), tw, th, "Pazarlama kanalı", "Açık kaynak küresel huni",
     "Apache-2.0 çekirdek GitHub/Show HN üzerinden dünya çapında geliştiriciye ulaşır, ücretli katmana dönüşür (Infisical/Authentik modeli).")
tile(s, Inches(7.0), Inches(4.35), tw, th, "Ölçeklenebilirlik", "Marjinal maliyet ≈ 0",
     "Bir yazılım kopyasını 1 veya 1.000 ülkeye satmanın üretim maliyeti aynıdır. Yüksek brüt marjlı, hızlı ölçeklenen ihracat kalemi.")
footer(s, 3)

# ---------- S4 section 2 product ----------
s = prs.slides.add_slide(BLANK); bg(s)
kicker_title(s, "BÖLÜM 2", "Sunulan ürün ve hizmet", kcolor=GOLD)
cw = Inches(2.95); ch = Inches(2.5); cy = Inches(1.85)
tile(s, Inches(0.65), cy, cw, ch, "IAM", "Kimlik & Erişim",
     "OIDC/OAuth2, SAML 2.0 IdP, SSO, MFA (TOTP, passkey, push, OTP), SCIM 2.0, LDAP/AD/Azure AD senkron, risk-tabanlı + step-up.")
tile(s, Inches(3.75), cy, cw, ch, "IGA", "Yönetişim",
     "Erişim gözden geçirme, SoD, çok-adımlı onay, JIT yükseltme, HR-tetikli Joiner/Mover/Leaver, RBAC/ABAC/OPA.")
tile(s, Inches(6.85), cy, cw, ch, "PAM", "Ayrıcalıklı Erişim",
     "Şifreli kimlik kasası + rotasyon, brokerlı SSH/RDP/VNC, şifreli oturum kaydı, ayrıcalık grafiği, risk-tabanlı askıya alma, dört-göz oturum.")
tile(s, Inches(9.95), cy, cw, ch, "ZTNA", "Sıfır-Güven Ağ",
     "OpenZiti overlay ile karanlık servisler (açık port yok), clientless erişim, masaüstü/mobil ajan + posture, reconciler.")
bullet_box(s, Inches(0.65), Inches(4.6), Inches(6.0), Inches(2.0),
           "Hizmetler", OK,
           ["Kurulum & yapılandırma, yönetilen dağıtım",
            "SLA + öncelikli destek abonelikleri",
            "Uyum paketleri (SOC 2/ISO/KVKK/DORA), premium konnektörler"],
           fill=RGBColor(0x0E,0x1E,0x18))
bullet_box(s, Inches(7.0), Inches(4.6), Inches(5.65), Inches(2.0),
           "Teknik temel", BRAND2,
           ["Go + PostgreSQL; CI-zorunlu çok-kiracılılık (FORCE RLS)",
            "Kurcalama-belirtir denetim (HMAC hash-zinciri), APISIX, mTLS"],
           fill=PANEL)
footer(s, 4)

# ---------- S5 moat quote ----------
s = prs.slides.add_slide(BLANK); bg(s)
txt(s, Inches(0.65), Inches(0.6), Inches(12), Inches(0.4),
    [[("REKABET FARKI", 14, BRAND2, True)]])
txt(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(3.0),
    [[("Tek bir kill-switch komutu; kullanıcının ", 28, WHITE, True),
      ("token'larını, oturumlarını, kasa check-out'larını, canlı ayrıcalıklı oturumlarını VE ağ devrelerini", 28, BRAND2, True),
      (" saniyeler içinde keser — çok-satıcılı entegrasyon projesi olarak değil, ", 28, WHITE, True),
      ("tek işlemle.", 28, BRAND2, True)]])
txt(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.0),
    [[("2026 sektör konsolidasyonu (PANW × CyberArk $25 milyar, Okta × Axiom, Delinea × StrongDM) "
       "bu birleşik ürünü M&A ile kurmaya çalışıyor. OpenIDX zaten tek kod tabanında birleşik.", 14, MUTED, False)]])
footer(s, 5)

# ---------- S6 section 3 target markets ----------
s = prs.slides.add_slide(BLANK); bg(s)
kicker_title(s, "BÖLÜM 3", "Hedef pazarlar", kcolor=GOLD)
txt(s, Inches(0.65), Inches(1.75), Inches(12), Inches(0.6),
    [[("Strateji: regülasyonun self-host'u zorunlu kıldığı pazarlardan başla (en kolay kazanım), "
       "sonra maliyet-hassas açık-kaynak-öncelikli kurumlara genişle.", 15, MUTED, False)]])
rows = [
    ("1 · Yurt içi çıpa", "Türkiye kamu & regüle sektör",
     "2019/12 Genelge + BDDK/KVKK yurt dışı SaaS IdP'yi yasaklıyor; yerli self-host tek uyumlu yol. İhracat için vitrin."),
    ("2 · İhracat", "AB — NIS2 / DORA orta ölçek",
     "DORA pratikte PAM zorunluluğu; NIS2 kimlik+ayrıcalıklı erişim+denetim ister. Egemen self-host tek-satıcı boşluğu."),
    ("2 · İhracat", "Körfez & MENA kamu/finans",
     "Veri yerelleştirme + egemenlik öncelikli bölge; self-host güvenlik yığınına güçlü talep."),
    ("3 · İhracat", "Küresel MSP'ler",
     "White-label çok-kiracılı ZTNA açık kaynakta yok; FORCE-RLS + ölçüm/geri-faturalama tam oturur."),
    ("3 · İhracat", "OSS-öncelikli mühendislik kurumları",
     "Auth0-yenileme ve Teleport-lisans mültecileri; GitHub/Show HN hunisi."),
]
ty = Inches(2.5); rh = Inches(0.86)
# header
txt(s, Inches(0.65), ty, Inches(2.3), Inches(0.3), [[("ÖNCELİK", 12, WHITE, True)]])
txt(s, Inches(3.0), ty, Inches(3.4), Inches(0.3), [[("PAZAR", 12, WHITE, True)]])
txt(s, Inches(6.5), ty, Inches(6.2), Inches(0.3), [[("NEDEN OPENIDX KAZANIR", 12, WHITE, True)]])
for i,(pri,mkt,why) in enumerate(rows):
    yy = ty + Inches(0.4) + rh*i
    txt(s, Inches(0.65), yy, Inches(2.3), rh, [[(pri, 13, BRAND2, True)]])
    txt(s, Inches(3.0), yy, Inches(3.4), rh, [[(mkt, 13, WHITE, True)]])
    txt(s, Inches(6.5), yy, Inches(6.2), rh, [[(why, 12, MUTED, False)]])
footer(s, 6)

# ---------- S7 pricing ----------
s = prs.slides.add_slide(BLANK); bg(s)
kicker_title(s, "Pazar & Fiyat Konumu", "Neden fiyatla kazanır")
bullet_box(s, Inches(0.65), Inches(1.9), Inches(6.0), Inches(2.55),
           "Alıcının acısı: fiyat", BAD,
           ["Okta Essentials ~$204/kullanıcı/yıl; governance ek SKU",
            "Auth0: 1.67× kullanıcı → belgelenmiş 15× fatura",
            "CyberArk PAM $1.8K–12K/kullanıcı/yıl, teklif-bazlı",
            "Teleport: kaynak-başı metering ölçeklenmeyi cezalandırır"],
           fill=RGBColor(0x22,0x0E,0x10))
bullet_box(s, Inches(7.0), Inches(1.9), Inches(5.65), Inches(2.55),
           "OpenIDX yanıtı", OK,
           ["Enterprise ~$4–6/kullanıcı/ay, sabit maliyet",
            "Yığılmış rakip $33–65 → %70–80 tasarruf",
            "Community sonsuza dek Apache-2.0 (huni)",
            "Kilitlenme yok, veri egemenliği"],
           fill=RGBColor(0x0E,0x1E,0x18))
pw = Inches(2.95); ph = Inches(1.85); py = Inches(4.75)
def ptier(l, h, price, x, hl=False):
    box(s, l, py, pw, ph, fill=PANEL, line=(BRAND2 if hl else LINE), line_w=(2 if hl else 1))
    txt(s, l+Inches(0.18), py+Inches(0.16), pw-Inches(0.36), ph-Inches(0.3),
        [[(h, 14, BRAND2, True)],[(price, 18, WHITE, True)],[(x, 11, MUTED, False)]], sp_after=3)
ptier(Inches(0.65), "Community", "Ücretsiz", "Tüm çekirdek, self-host")
ptier(Inches(3.75), "Enterprise", "$4–6 /kul/ay", "SLA, uyum, premium konnektör", hl=True)
ptier(Inches(6.85), "MSP", "Kiracı + ölçüm", "White-label çok-kiracı")
ptier(Inches(9.95), "Kamu/Savunma", "Proje-bazlı", "Uyumlu self-host, yerli malı, FIPS")
footer(s, 7)

# ---------- S8 section 4 stage ----------
s = prs.slides.add_slide(BLANK); bg(s)
kicker_title(s, "BÖLÜM 4", "Mevcut gelişim aşaması", kcolor=GOLD)
b = box(s, Inches(0.65), Inches(1.75), Inches(8.5), Inches(0.5), fill=PANEL, line=GOLD)
txt(s, Inches(0.85), Inches(1.85), Inches(8.2), Inches(0.4),
    [[("Olgunluk: MVP–Erken GA · çekirdek gerçek ve büyük ölçüde üretim kalitesinde", 12.5, GOLD, True)]])
bullet_box(s, Inches(0.65), Inches(2.5), Inches(6.0), Inches(2.35),
           "Hazır & çalışır (kod-doğrulanmış)", OK,
           ["Dört sütun canlı; tek Postgres referans ortam",
            "PAM derinliği: ayrıcalık grafiği, risk-askıya alma, dört-göz, şifreli kayıt — canlı doğrulandı",
            "Native erişim: kısa-ömürlü SSH sertifikası CLI (openidx-connect)",
            "CI-zorunlu çok-kiracılılık, kurcalama-belirtir denetim"],
           fill=RGBColor(0x0E,0x1E,0x18))
bullet_box(s, Inches(7.0), Inches(2.5), Inches(5.65), Inches(2.35),
           "Devam eden yol haritası", BRAND2,
           ["DB/K8s/cloud oturum brokerları",
            "ChatOps onayları, AI-agent erişim derinliği",
            "Uyum: Pentest → ISO 27001 → SOC 2 Type II",
            "CRA (AB Siber Dayanıklılık Yasası) hazırlığı"],
           fill=PANEL)
kw2 = Inches(2.95); kh2 = Inches(1.55); ky2 = Inches(5.1)
kpi(s, Inches(0.65), ky2, kw2, kh2, "Canlı", "Üretim benzeri referans ortam")
kpi(s, Inches(3.75), ky2, kw2, kh2, "Pilot", "Ücretli pilot & tasarım-ortağı için hazır")
kpi(s, Inches(6.85), ky2, kw2, kh2, "Açık", "Apache-2.0 çekirdek")
kpi(s, Inches(9.95), ky2, kw2, kh2, "Yerli", "Tamamen yerli geliştirme")
footer(s, 8)

# ---------- S9 section 5 growth ----------
s = prs.slides.add_slide(BLANK); bg(s)
kicker_title(s, "BÖLÜM 5", "Global büyüme hedefleri", kcolor=GOLD)
# flow
from pptx.enum.shapes import MSO_SHAPE
nodes = ["Yurt içi vitrin", "AB (NIS2/DORA)", "Körfez/MENA", "Küresel MSP + OSS"]
nx = Inches(0.65); nw = Inches(2.7); nh = Inches(0.8); nyy = Inches(1.8)
for i,n in enumerate(nodes):
    l = nx + (nw+Inches(0.32))*i
    box(s, l, nyy, nw, nh, fill=PANEL, line=LINE)
    txt(s, l, nyy+Inches(0.22), nw, Inches(0.4), [[(n, 14, WHITE, True)]], align=PP_ALIGN.CENTER)
    if i < 3:
        txt(s, l+nw, nyy+Inches(0.15), Inches(0.32), Inches(0.5), [[("→", 20, BRAND2, True)]], align=PP_ALIGN.CENTER)
tw2 = Inches(3.9); th2 = Inches(1.7)
tile(s, Inches(0.65), Inches(3.0), tw2, th2, "Hedef", "Yıl-2: $1–2M ARR",
     "30–60 ihracat + yurt içi müşteri. Zitadel/NetBird/Infisical su hattı.")
tile(s, Inches(4.72), Inches(3.0), tw2, th2, "Kanal", "Açık kaynak huni + iş ortakları",
     "GitHub/Show HN küresel görünürlük; yerel entegratörler; MSP white-label.")
tile(s, Inches(8.79), Inches(3.0), tw2, th2, "Güven zarfı", "Uluslararası sertifikalar",
     "ISO 27001, SOC 2 Type II, CRA hazırlığı — kurumsal satış önkoşulu.")
bullet_box(s, Inches(0.65), Inches(4.95), Inches(6.0), Inches(1.9),
           "E-ihracat desteğinin katkısı", OK,
           ["Yurt dışı dijital pazarlama, çok-dilli içerik, SEO/topluluk",
            "Uluslararası sertifikasyon ve hedef-pazar uyum danışmanlığı",
            "Yurt dışı fuar/etkinlik ve iş ortağı ağı geliştirme"],
           fill=RGBColor(0x0E,0x1E,0x18))
bullet_box(s, Inches(7.0), Inches(4.95), Inches(5.65), Inches(1.9),
           "Ölçülebilir çıktı", BRAND2,
           ["Yurt dışı ödemeli müşteri sayısı ve döviz geliri",
            "Aktif ülke sayısı ve topluluk (GitHub) büyüklüğü"],
           fill=PANEL)
footer(s, 9)

# ---------- S10 closing ----------
s = prs.slides.add_slide(BLANK); bg(s)
txt(s, Inches(0.65), Inches(0.7), Inches(6), Inches(0.5),
    [[("◇  Open", 22, WHITE, True), ("IDX", 22, BRAND2, True)]])
txt(s, Inches(0.65), Inches(1.9), Inches(12), Inches(1.4),
    [[("Yerli üretim, self-host, dört-sütunlu siber güvenlik", 34, WHITE, True)],
     [("platformu — dijital ihracata hazır", 34, BRAND2, True)]], sp_after=2)
bullet_box(s, Inches(0.65), Inches(3.8), Inches(9.5), Inches(2.0),
           "Özet", OK,
           ["Dört ürün tek platformda, %70–80 daha uygun maliyet",
            "Regülasyonun self-host'u zorunlu kıldığı küresel pazarlar hedefte",
            "MVP–Erken GA; ücretli pilot ve global büyümeye hazır"],
           fill=RGBColor(0x0E,0x1E,0x18))
txt(s, Inches(0.65), Inches(6.0), Inches(12), Inches(0.5),
    [[("İletişim · ", 16, INK, False), ("openidx.tdv.org", 16, WHITE, True),
      (" · github.com/mhmtgngr/openidx", 16, MUTED, False)]])
footer(s, 10)

out = "OpenIDX-E-Ihracat-Sunum-TR.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
